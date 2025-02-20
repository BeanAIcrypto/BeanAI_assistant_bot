import zipfile
import rarfile
import py7zr
import tarfile
import os
import PyPDF2
from docx import Document
from odf.opendocument import load
from odf.text import P
import pandas as pd
from pptx import Presentation
import logging
from bs4 import BeautifulSoup
import email
import mimetypes
import tempfile

from typing import List, Optional

logger = logging.getLogger(__name__)


def extract_text_from_pdf(file_path: str) -> str:
    """
    Извлекает текст из PDF файла.

    Args:
        file_path (str): Путь к PDF файлу.

    Returns:
        str: Извлеченный текст из всех страниц PDF.

    Raises:
        FileNotFoundError: Файл PDF не найден.
        ValueError: Ошибка чтения PDF файла.
    """
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            text = "".join(page.extract_text() or "" for page in reader.pages)
        return text
    except FileNotFoundError as e:
        logger.error(f"Файл PDF не найден: {file_path}")
        raise e
    except PyPDF2.errors.PdfReadError as e:
        logger.error(f"Ошибка чтения PDF файла {file_path}: {e}")
        raise ValueError("Не удалось прочитать PDF файл.")
    except Exception as e:
        logger.error(f"Неизвестная ошибка при обработке PDF {file_path}: {e}")
        raise ValueError("Произошла ошибка при извлечении текста из PDF.")


def extract_text_from_txt(file_path: str) -> str:
    """
    Извлекает текст из текстового файла.

    Args:
        file_path (str): Путь к текстовому файлу (TXT).

    Returns:
        str: Извлеченный текст из файла.

    Raises:
        FileNotFoundError: Файл не найден.
        UnicodeDecodeError: Ошибка декодирования файла.
        IOError: Ошибка ввода/вывода при чтении файла.
    """
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except FileNotFoundError as e:
        logger.error(f"Файл TXT не найден: {file_path}")
        raise e
    except UnicodeDecodeError as e:
        logger.error(f"Ошибка декодирования TXT файла {file_path}: {e}")
        raise UnicodeDecodeError(
            f"Ошибка декодирования текста в файле {file_path}."
        )
    except IOError as e:
        logger.error(
            f"Ошибка ввода/вывода при чтении TXT файла {file_path}: {e}"
        )
        raise IOError(f"Не удалось прочитать файл {file_path}.")
    except Exception as e:
        logger.error(
            f"Неизвестная ошибка при обработке TXT файла {file_path}: {e}"
        )
        raise Exception("Произошла ошибка при извлечении текста из TXT файла.")


def extract_text_from_docx(file_path: str) -> str:
    """
    Извлекает текст из DOCX файла.

    Args:
        file_path (str): Путь к файлу формата DOCX.

    Returns:
        str: Извлеченный текст из файла.

    Raises:
        FileNotFoundError: Файл не найден.
        ValueError: Файл имеет некорректный формат.
        IOError: Ошибка ввода/вывода при чтении файла.
    """
    try:
        _, tipy_doc = os.path.splitext(file_path)
        if tipy_doc == ".doc":
            logger.error(f"Формат .doc не поддерживается: {file_path}")
            return
        doc: Document = Document(file_path)
        paragraphs: List[str] = [para.text for para in doc.paragraphs]
        return "\n".join(paragraphs)
    except FileNotFoundError as e:
        logger.error(f"Файл DOCX не найден: {file_path}")
        raise FileNotFoundError(f"Файл {file_path} не найден.")
    except ValueError as e:
        logger.error(
            f"Файл DOCX поврежден или имеет неверный формат: {file_path}"
        )
        raise ValueError(
            f"Файл {file_path} поврежден или не является корректным DOCX файлом."
        )
    except IOError as e:
        logger.error(
            f"Ошибка ввода/вывода при чтении DOCX файла {file_path}: {e}"
        )
        raise IOError(f"Не удалось прочитать файл {file_path}.")
    except Exception as e:
        logger.error(
            f"Неизвестная ошибка при обработке DOCX файла {file_path}: {e}"
        )
        raise Exception(
            f"Произошла ошибка при извлечении текста из файла {file_path}."
        )


def extract_text_from_odt(file_path: str) -> str:
    """
    Извлекает текст из ODT файла.

    Args:
        file_path (str): Путь к файлу формата ODT.

    Returns:
        str: Извлеченный текст из файла.

    Raises:
        FileNotFoundError: Файл не найден.
        ValueError: Файл имеет некорректный формат.
        IOError: Ошибка ввода/вывода при чтении файла.
    """
    try:
        doc = load(file_path)
        texts: List[str] = [
            text.firstChild.data
            for text in doc.getElementsByType(P)
            if text.firstChild is not None
        ]
        return "\n".join(texts)
    except FileNotFoundError as e:
        logger.error(f"Файл ODT не найден: {file_path}")
        raise FileNotFoundError(f"Файл {file_path} не найден.")
    except ValueError as e:
        logger.error(f"Некорректный формат файла ODT: {file_path}")
        raise ValueError(
            f"Файл {file_path} поврежден или не является корректным ODT файлом."
        )
    except IOError as e:
        logger.error(
            f"Ошибка ввода/вывода при чтении ODT файла {file_path}: {e}"
        )
        raise IOError(f"Не удалось прочитать файл {file_path}.")
    except Exception as e:
        logger.error(
            f"Неизвестная ошибка при обработке ODT файла {file_path}: {e}"
        )
        raise Exception(
            f"Произошла ошибка при извлечении текста из файла {file_path}."
        )


def extract_text_from_excel(file_path: str) -> Optional[str]:
    """
    Извлекает текст из файла Excel (XLS, XLSX, ODS).

    Args:
        file_path (str): Путь к файлу формата Excel.

    Returns:
        Optional[str]: Текст из файла в формате Markdown или None в случае ошибки.

    Raises:
        FileNotFoundError: Файл не найден.
        ValueError: Файл пустой или некорректный.
        IOError: Ошибка ввода/вывода.
    """
    try:
        df = pd.read_excel(file_path)
        df = df.loc[:, ~df.columns.str.contains("^Unnamed")]
        df = df.dropna(how="all").dropna(axis=1, how="all")

        if df.empty:
            raise ValueError(f"Файл {file_path} пустой.")

        return df.to_markdown(index=False)
    except FileNotFoundError as e:
        logger.error(f"Файл не найден: {file_path}")
        raise FileNotFoundError(f"Файл {file_path} не найден.") from e
    except ValueError as e:
        logger.error(f"Файл пустой или некорректный: {file_path}")
        raise ValueError(f"Файл {file_path} пустой или некорректный.") from e
    except IOError as e:
        logger.error(f"Ошибка ввода/вывода: {file_path}")
        raise IOError(f"Ошибка ввода/вывода при чтении {file_path}.") from e
    except Exception as e:
        logger.error(f"Ошибка обработки файла: {file_path}")
        raise Exception(f"Ошибка при обработке файла {file_path}.") from e


def extract_text_from_presentation(file_path: str) -> Optional[str]:
    """
    Извлекает текст из файла PowerPoint (PPT, PPTX, ODP).

    Args:
        file_path (str): Путь к файлу презентации.

    Returns:
        Optional[str]: Извлеченный текст из презентации или None в случае ошибки.

    Raises:
        FileNotFoundError: Файл не найден.
        ValueError: Файл пустой или некорректный.
        IOError: Ошибка ввода/вывода.
    """
    try:
        presentation = Presentation(file_path)
        text = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text.strip())

        if not text:
            raise ValueError(f"Файл {file_path} не содержит текста.")

        return "\n".join(text)
    except FileNotFoundError as e:
        logger.error(f"Файл не найден: {file_path}")
        raise FileNotFoundError(f"Файл {file_path} не найден.") from e
    except ValueError as e:
        logger.error(f"Файл пустой или не содержит текста: {file_path}")
        raise ValueError(
            f"Файл {file_path} пустой или не содержит текста."
        ) from e
    except IOError as e:
        logger.error(f"Ошибка ввода/вывода при обработке файла: {file_path}")
        raise IOError(f"Ошибка ввода/вывода при чтении {file_path}.") from e
    except Exception as e:
        logger.error(f"Ошибка обработки презентации: {file_path}")
        raise Exception(f"Ошибка при обработке файла {file_path}.") from e


def extract_text_from_zip(file_path: str) -> Optional[str]:
    """
    Извлекает текст из ZIP-архива с учетом MIME-типа.

    Args:
        file_path (str): Путь к ZIP-архиву.

    Returns:
        Optional[str]: Извлеченный текст из поддерживаемых файлов или None в случае ошибки.

    Raises:
        FileNotFoundError: ZIP-файл не найден.
        ValueError: ZIP-файл пустой или некорректный.
        IOError: Ошибка ввода/вывода при чтении файлов из архива.
    """
    try:
        extracted_text: List[str] = []

        with tempfile.TemporaryDirectory() as tmpdir:
            with zipfile.ZipFile(file_path, "r") as archive:
                if not archive.namelist():
                    raise ValueError(f"Архив {file_path} пустой.")

                for file in archive.namelist():
                    mime_type = (
                        mimetypes.guess_type(file)[0]
                        or "application/octet-stream"
                    )

                    if mime_type in text_extraction_from_a_document:
                        extracted_file_path = archive.extract(
                            file, path=tmpdir
                        )
                        text_extraction_function = (
                            text_extraction_from_a_document[mime_type]
                        )
                        text_document = text_extraction_function(
                            extracted_file_path
                        )

                        if text_document:
                            extracted_text.append(text_document)
                        else:
                            logger.error(
                                f"Не удалось извлечь текст из файла в ZIP: {file}"
                            )
                    else:
                        logger.error(
                            f"Файл в ZIP имеет неподдерживаемый формат: {file} (MIME-тип: {mime_type})"
                        )

        if not extracted_text:
            raise ValueError(
                f"Архив {file_path} не содержит файлов с поддерживаемым текстовым форматом."
            )

        return "\n".join(extracted_text)
    except FileNotFoundError as e:
        logger.error(f"ZIP-файл не найден: {file_path}")
        raise FileNotFoundError(f"Файл {file_path} не найден.") from e
    except ValueError as e:
        logger.error(f"Ошибка обработки содержимого ZIP-файла: {file_path}")
        raise ValueError(
            f"Ошибка обработки содержимого файла {file_path}."
        ) from e
    except IOError as e:
        logger.error(f"Ошибка ввода/вывода при чтении ZIP-файла: {file_path}")
        raise IOError(f"Ошибка ввода/вывода при чтении {file_path}.") from e
    except Exception as e:
        logger.error(
            f"Неизвестная ошибка при обработке ZIP-файла: {file_path}"
        )
        raise Exception(
            f"Произошла ошибка при обработке файла {file_path}."
        ) from e


def extract_text_from_rar(file_path: str) -> Optional[str]:
    return "Данный формат не поддерживается. Сконвертируйте, пожалуйста в .zip"


def extract_text_from_7z(file_path: str) -> Optional[str]:
    """
    Извлекает текст из файлов в 7Z-архиве с учетом MIME-типа.

    Args:
        file_path (str): Путь к 7Z-архиву.

    Returns:
        Optional[str]: Извлеченный текст из поддерживаемых файлов или None в случае ошибки.

    Raises:
        FileNotFoundError: Если файл архива не найден.
        ValueError: Если архив пуст или не содержит поддерживаемых файлов.
        py7zr.SevenZipFileError: Если файл не является корректным 7Z-архивом.
        IOError: Ошибка ввода-вывода при обработке файла.
    """
    try:
        extracted_text = []

        with py7zr.SevenZipFile(file_path, mode="r") as archive:
            if not archive.getnames():
                raise ValueError(
                    f"7Z-архив '{file_path}' пуст или некорректен."
                )

            archive.extractall(path="tmp")
            for root, _, files in os.walk("tmp"):
                for file in files:
                    mime_type = (
                        mimetypes.guess_type(file)[0]
                        or "application/octet-stream"
                    )
                    file_path_extracted = os.path.join(root, file)

                    if mime_type in text_extraction_from_a_document:
                        text_extraction_function = (
                            text_extraction_from_a_document[mime_type]
                        )
                        with open(
                            file_path_extracted,
                            "r",
                            encoding="utf-8",
                            errors="ignore",
                        ) as f:
                            text_content = text_extraction_function(
                                file_path_extracted
                            )
                            if text_content:
                                extracted_text.append(text_content)
                            else:
                                logger.error(
                                    f"Не удалось извлечь текст из файла: {file}"
                                )
                    else:
                        logger.error(
                            f"Неподдерживаемый формат файла в 7Z архиве: {file} (MIME-тип: {mime_type})"
                        )

        if not extracted_text:
            raise ValueError(
                f"7Z-архив '{file_path}' не содержит поддерживаемых файлов."
            )

        return "\n".join(extracted_text)

    except FileNotFoundError as e:
        logger.error(f"Файл 7Z не найден: {file_path}")
        raise FileNotFoundError(f"7Z-файл '{file_path}' не найден.") from e
    except py7zr.SevenZipFileError as e:
        logger.error(f"Некорректный 7Z-архив: {file_path}")
        raise py7zr.SevenZipFileError(
            f"Файл '{file_path}' не является корректным 7Z-архивом."
        ) from e
    except IOError as e:
        logger.error(f"Ошибка ввода-вывода при обработке 7Z-архива: {e}")
        raise IOError(
            f"Ошибка ввода-вывода при обработке файла '{file_path}'."
        ) from e
    except Exception as e:
        logger.error(f"Произошла ошибка при обработке 7Z-архива: {e}")
        raise Exception(
            f"Ошибка при обработке 7Z-архива '{file_path}': {str(e)}"
        ) from e


def extract_text_from_tar(file_path: str) -> Optional[str]:
    """
    Извлекает текст из файлов в TAR-архиве с учетом MIME-типа.

    Args:
        file_path (str): Путь к TAR-архиву.

    Returns:
        Optional[str]: Извлеченный текст из поддерживаемых файлов или None в случае ошибки.

    Raises:
        FileNotFoundError: Если TAR-архив не найден.
        ValueError: Если архив пуст или не содержит поддерживаемых файлов.
        tarfile.ReadError: Если файл не является корректным TAR-архивом.
        IOError: Ошибка ввода-вывода при обработке файла.
    """
    try:
        extracted_text = []

        with tarfile.open(file_path, "r") as archive:
            if not archive.getmembers():
                raise ValueError(
                    f"TAR-архив '{file_path}' пуст или некорректен."
                )

            for member in archive.getmembers():
                if member.isfile():
                    mime_type = (
                        mimetypes.guess_type(member.name)[0]
                        or "application/octet-stream"
                    )

                    if mime_type in text_extraction_from_a_document:
                        with archive.extractfile(member) as f:
                            if f is not None:
                                text_extraction_function = (
                                    text_extraction_from_a_document[mime_type]
                                )
                                text_content = text_extraction_function(f)
                                if text_content:
                                    extracted_text.append(text_content)
                                else:
                                    logger.error(
                                        f"Не удалось извлечь текст из файла: {member.name}"
                                    )
                    else:
                        logger.error(
                            f"TAR-архив содержит неподдерживаемый формат: {member.name} (MIME-тип: {mime_type})"
                        )

        if not extracted_text:
            raise ValueError(
                f"TAR-архив '{file_path}' не содержит поддерживаемых файлов."
            )

        return "\n".join(extracted_text)

    except FileNotFoundError as e:
        logger.error(f"Файл TAR не найден: {file_path}")
        raise FileNotFoundError(f"TAR-файл '{file_path}' не найден.") from e
    except tarfile.ReadError as e:
        logger.error(f"Некорректный TAR-архив: {file_path}")
        raise tarfile.ReadError(
            f"Файл '{file_path}' не является корректным TAR-архивом."
        ) from e
    except IOError as e:
        logger.error(f"Ошибка ввода-вывода при обработке TAR-архива: {e}")
        raise IOError(
            f"Ошибка ввода-вывода при обработке файла '{file_path}'."
        ) from e
    except Exception as e:
        logger.error(f"Произошла ошибка при обработке TAR-архива: {e}")
        raise Exception(
            f"Ошибка при обработке TAR-архива '{file_path}': {str(e)}"
        ) from e


def extract_text_from_mhtml(file_path: str) -> Optional[str]:
    """
    Извлекает текст из MHTML-файла, включая текстовые и HTML-содержимые части.

    Args:
        file_path (str): Путь к MHTML-файлу.

    Returns:
        Optional[str]: Извлеченный текст или None в случае ошибки.

    Raises:
        FileNotFoundError: Если файл не найден.
        ValueError: Если файл пуст или не содержит текстовых данных.
        UnicodeDecodeError: Если не удается декодировать содержимое файла.
    """
    try:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Файл '{file_path}' не найден.")

        extracted_text = []

        with open(file_path, "rb") as file:
            msg = email.message_from_binary_file(file)

            for part in msg.walk():
                content_type = part.get_content_type()
                if content_type in ["text/plain", "text/html"]:
                    payload = part.get_payload(decode=True)
                    if payload:
                        charset = part.get_content_charset() or "utf-8"
                        try:
                            text = payload.decode(charset, errors="ignore")
                        except UnicodeDecodeError as decode_error:
                            logger.error(
                                f"Ошибка декодирования содержимого файла {file_path}: {decode_error}"
                            )
                            raise UnicodeDecodeError(
                                f"Ошибка декодирования файла '{file_path}'."
                            )

                        if content_type == "text/html":
                            soup = BeautifulSoup(text, "html.parser")
                            text = soup.get_text(separator="\n")

                        extracted_text.append(text.strip())

        if extracted_text:
            return "\n".join(extracted_text)
        else:
            raise ValueError(
                f"MHTML-файл '{file_path}' не содержит текстового содержимого."
            )

    except FileNotFoundError as e:
        logger.error(f"Файл не найден: {e}")
        raise FileNotFoundError(f"Файл '{file_path}' не найден.") from e
    except UnicodeDecodeError as e:
        logger.error(f"Ошибка декодирования: {e}")
        raise UnicodeDecodeError(
            f"Не удалось декодировать содержимое файла '{file_path}'."
        ) from e
    except ValueError as e:
        logger.error(f"Файл пуст или некорректный: {e}")
        raise ValueError(
            f"MHTML-файл '{file_path}' не содержит текстового содержимого."
        ) from e
    except Exception as e:
        logger.error(f"Ошибка при обработке MHTML-файла: {e}")
        raise Exception(
            f"Произошла ошибка при обработке файла '{file_path}': {str(e)}"
        ) from e


def extract_text_from_markdown(file_path: str) -> Optional[str]:
    """
    Извлекает текст из Markdown-файла.

    Args:
        file_path (str): Путь к Markdown-файлу.

    Returns:
        Optional[str]: Извлеченный текст из файла.

    Raises:
        FileNotFoundError: Если файл не найден.
        UnicodeDecodeError: Если не удается декодировать содержимое файла.
        ValueError: Если файл пуст или содержит некорректное содержимое.
    """
    try:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Файл '{file_path}' не найден.")

        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
            if not content.strip():
                raise ValueError(f"Файл '{file_path}' пуст или некорректен.")

        logger.info(f"Текст успешно извлечен из Markdown-файла: {file_path}")
        return content

    except FileNotFoundError as e:
        logger.error(f"Файл не найден: {e}")
        raise FileNotFoundError(f"Файл '{file_path}' не найден.") from e
    except UnicodeDecodeError as e:
        logger.error(f"Ошибка декодирования содержимого файла: {e}")
        raise UnicodeDecodeError(
            f"Не удалось декодировать содержимое файла '{file_path}'."
        ) from e
    except ValueError as e:
        logger.error(f"Ошибка чтения содержимого файла: {e}")
        raise ValueError(f"Файл '{file_path}' пуст или некорректен.") from e
    except Exception as e:
        logger.error(f"Произошла ошибка при обработке Markdown-файла: {e}")
        raise Exception(
            f"Ошибка при обработке файла '{file_path}': {str(e)}"
        ) from e


text_extraction_from_a_document = {
    "application/pdf": extract_text_from_pdf,
    "text/plain": extract_text_from_txt,
    "application/msword": extract_text_from_docx,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_text_from_docx,
    "application/vnd.oasis.opendocument.text": extract_text_from_odt,
    "application/vnd.ms-excel": extract_text_from_excel,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": extract_text_from_excel,
    "application/vnd.oasis.opendocument.spreadsheet": extract_text_from_excel,
    "application/vnd.ms-powerpoint": extract_text_from_presentation,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_text_from_presentation,
    "application/vnd.oasis.opendocument.presentation": extract_text_from_presentation,
    "application/zip": extract_text_from_zip,
    "application/x-rar-compressed": extract_text_from_rar,
    "application/x-7z-compressed": extract_text_from_7z,
    "application/x-tar": extract_text_from_tar,
    "application/gzip": extract_text_from_tar,
    "application/x-bzip2": extract_text_from_tar,
    "text/markdown": extract_text_from_markdown,
    "application/x-mimearchive": extract_text_from_mhtml,
    "text/html": extract_text_from_mhtml,
}
